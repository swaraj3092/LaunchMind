import sys
import os

# Install python-pptx if not present
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Sakura Noir Palette Colors
    BG_COLOR = RGBColor(12, 8, 13)       # Deep Plum background
    ROSE_COLOR = RGBColor(255, 45, 120)  # Neon Rose titles
    VIOLET_COLOR = RGBColor(192, 132, 252) # Soft Violet accents
    WHITE_COLOR = RGBColor(255, 240, 245) # Sakura White body text
    GRAY_COLOR = RGBColor(155, 126, 160)  # Muted secondary text
    
    slides = [
        {
            "num": "01",
            "title": "LaunchMind",
            "tagline": "Curing the Builder's Blindspot. Validate before you build.",
            "bullets": [
                "AI-Powered Startup Validation & Interactive VC Pitch Engine",
                "Built by: Swaraj Kumar Behera & Prajakta Kuila (Team: Dynamic Duo)",
                "For: InnovaHack Chapter 1 — Generative AI / Startup Innovation Track",
                "Web App URL: http://localhost:5173 / GitHub: swaraj3092/LaunchMind"
            ],
            "note": "We are introducing LaunchMind, the AI co-founder that helps you validate ideas before writing code."
        },
        {
            "num": "02",
            "title": "The Problem",
            "tagline": "The Startup Graveyard — Building blind.",
            "bullets": [
                "Echo Chambers: Founders validate ideas with supportive friends instead of market realities.",
                "Premature Optimization: Jumping straight into code before validating customer pain points.",
                "Consultancy Barrier: Getting expert startup advice or VC feedback is expensive and out of reach.",
                "Resource Loss: Billions of dollars and millions of hours of builder potential wasted annually."
            ],
            "note": "Technical founders are biased toward action. They write code first and validate later, creating solutions for non-existent problems."
        },
        {
            "num": "03",
            "title": "The Builder's Blindspot",
            "tagline": "35% of startups fail because of 'No Market Need'.",
            "bullets": [
                "Over 90% of tech startups fail within the first 3 years.",
                "CB Insights flags 'No Market Need' as the #1 cause of startup post-mortems (35%).",
                "Builders spend 6 months writing code before spending 6 minutes validating market realities.",
                "LaunchMind exists to bridge the gap between initial ideation and actual market validation."
            ],
            "note": "Wasting 6 months building something nobody wants is the single biggest risk in startup creation. We solve this Day 1."
        },
        {
            "num": "04",
            "title": "Our Solution",
            "tagline": "An always-on AI co-founder and Silicon Valley VC.",
            "bullets": [
                "Idea Interrogation: A personalized 5-question engine that grills the founder's assumptions.",
                "Live Web Intel: Real-time search scraping to find active, live competitors.",
                "Assumption Kill-Check: Surfaces the top 3 unvalidated risks that could kill the startup.",
                "Action Roadmap: A structured 30/60/90-day plan with an actionable Day 1 checklist."
            ],
            "note": "LaunchMind extracts assumptions, runs live web comparisons, and hands you an actionable roadmap in minutes."
        },
        {
            "num": "05",
            "title": "How It Works",
            "tagline": "The 4-step validation funnel.",
            "bullets": [
                "1. Raw Idea Input: The founder describes their vision in one simple sentence.",
                "2. AI Interrogation: Answer 5 targeted questions (Normal or Roast Mode 🔥).",
                "3. Plan Synthesis: AI creates roadmap, lists competitors, and calculates market size.",
                "4. Live VC Pitch: Defend the idea in the Pitch Room and get a Venture Readiness Score."
            ],
            "note": "Our workflow guides the founder from raw ideation, to competitive defense, to a formatted execution plan."
        },
        {
            "num": "06",
            "title": "The VC Pitch Room",
            "tagline": "Interactive VC pitch room with real-time score tracking.",
            "bullets": [
                "3 AI Personas: Marcus Vance (Skeptic), Elena Rostova (Growth), Dr. Aris Thorne (Deep Tech).",
                "Real-time Scoring: Real-time scoreboard updates (e.g. +10 Strong, -5 Red Flag) after each response.",
                "AI-driven Due Diligence: 4 rounds of sharp, semantic follow-up questions tailored to your inputs.",
                "Final Verdict: A blunt partner-level decision (Invest/Pass) with a complete feedback log."
            ],
            "note": "This is our key differentiator. Founders get grilled by simulated partners to test their market and tech defensibility."
        },
        {
            "num": "07",
            "title": "Competitive Advantage",
            "tagline": "LaunchMind vs. Traditional Alternatives.",
            "bullets": [
                "Real-Time Scraping: Unlike database tools, we run live search queries to catch active competitors.",
                "Interactive VC: Moving past static plans to an active, conversational partner simulator.",
                "Frictionless UX: Fast, glassmorphic React SPA designed for modern developers.",
                "Action-Oriented: Every plan centers around a realistic, immediate Day 1 validation task."
            ],
            "note": "LaunchMind is faster, live-connected, and conversational — moving past static templates to interactive validation."
        },
        {
            "num": "08",
            "title": "Technical Architecture",
            "tagline": "Built for speed, safety, and scale.",
            "bullets": [
                "Frontend: React SPA using Vite, TypeScript, TailwindCSS, and Framer Motion animations.",
                "Backend: FastAPI (Python) asynchronous REST API endpoints.",
                "AI Pipeline: Google Gemini 2.5 Flash with structured JSON output templates.",
                "Integrations: DuckDuckGo live search api wrapper for competitor mining."
            ],
            "note": "Our stateless backend scales horizontally, using asynchronous parallel requests to Gemini and search to keep response times low."
        },
        {
            "num": "09",
            "title": "Traction & Roadmap",
            "tagline": "Our execution plan.",
            "bullets": [
                "Phase 1 (Now): Core validation engine and VC Pitch Room MVP completed.",
                "Phase 2 (Q3): User authentication, project dashboard persistence, and progress tracking.",
                "Phase 3 (Q4): Voice-to-voice pitch capability and vector-based YC deck benchmark scanning.",
                "Target: 10,000 active validations and partnerships with university hackathon organizers."
            ],
            "note": "We have the core MVP ready. Next, we persist user sessions, add voice capability, and integrate historical YC benchmarking."
        },
        {
            "num": "10",
            "title": "The Ask",
            "tagline": "Partnering to scale startup validation.",
            "bullets": [
                "Stage: Working prototype deployed and open-source on GitHub.",
                "Incubator Partnerships: License to university networks, venture studios, and incubators.",
                "Hackathon Integrations: Position LaunchMind as the official pre-validation tool for hackathons.",
                "Vision: To become the global standard for pre-seed project validation."
            ],
            "note": "We're looking to integrate with hackathons and incubators to help founders test assumptions before spending code resources."
        },
        {
            "num": "11",
            "title": "Business Model",
            "tagline": "Scalable SaaS monetization strategy.",
            "bullets": [
                "Freemium Tier: 3 free startup idea validations per month for individual developers.",
                "Pro Subscription ($19/mo): Unlimited idea tests, full VC Room access, and advanced PDF exports.",
                "Enterprise Licensing: Dedicated validation dashboards for university hubs and accelerators.",
                "API License: Licensing validation APIs to developer platforms and startup application portals."
            ],
            "note": "We monetize through a premium subscription for power creators, and enterprise packages for incubators."
        },
        {
            "num": "12",
            "title": "Go-To-Market",
            "tagline": "Reaching builders at the source.",
            "bullets": [
                "Hackathon Partnerships: Standardize LaunchMind as the pre-submission check for hackathons.",
                "GitHub Action: Trigger a validation audit on repository creation via GitHub marketplace.",
                "Build In Public: Drive organic traffic via viral LinkedIn and Twitter posts showing VC roasts.",
                "Startup Hub Onboarding: Partner with incubator networks to use LaunchMind for cohort screening."
            ],
            "note": "Our GTM puts LaunchMind right where developers start — hackathons, repository creation, and early application portals."
        },
        {
            "num": "13",
            "title": "The Team",
            "tagline": "Team Dynamic Duo — Curing the Builder's Blindspot.",
            "bullets": [
                "Swaraj Kumar Behera — Fullstack Developer & AI Integration Specialist.",
                "Prajakta Kuila — Frontend Engineer & Lead UI/UX Designer.",
                "Our Goal: Making startup validation accessible, fast, and objective.",
                "GitHub Repo: https://github.com/swaraj3092/LaunchMind"
            ],
            "note": "We are Swaraj and Prajakta, team Dynamic Duo, committed to helping developers build products people actually want."
        }
    ]
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set background to dark plum
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Slide number indicator (Sakura theme accent)
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(2), Inches(0.8))
        tf_num = num_box.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = f"// SLIDE {slide_data['num']}"
        p_num.font.name = "Courier New"
        p_num.font.size = Pt(14)
        p_num.font.bold = True
        p_num.font.color.rgb = VIOLET_COLOR
        
        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = ROSE_COLOR
        
        # Tagline Box
        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.5))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = slide_data["tagline"]
        p_tag.font.name = "Arial"
        p_tag.font.size = Pt(18)
        p_tag.font.italic = True
        p_tag.font.color.rgb = VIOLET_COLOR
        
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(4.2))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for bullet_text in slide_data["bullets"]:
            p = tf_content.add_paragraph()
            p.text = "• " + bullet_text
            p.space_after = Pt(14)
            p.font.name = "Arial"
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE_COLOR
            
        # Add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data["note"]

    # Save presentation
    output_path = "LaunchMind_PitchDeck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
